"""
Creates WikiMed_Presentation.pptx in the current directory.

Upload to Google Drive (drive.google.com †' New †' File upload) and
Google will automatically open it in Google Slides for editing.

Usage:
    python create_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import pptx.oxml.ns as ns
from lxml import etree

# "" Palette """"""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""
BG      = RGBColor(0x08, 0x14, 0x1F)
SURFACE = RGBColor(0x0E, 0x1E, 0x2D)
BORDER  = RGBColor(0x1B, 0x33, 0x49)
TEXT    = RGBColor(0xDC, 0xF0, 0xFA)
MUTED   = RGBColor(0x6A, 0x9A, 0xB5)
DIM     = RGBColor(0x3D, 0x6A, 0x8A)
ACCENT  = RGBColor(0x00, 0xC5, 0xD8)
WARM    = RGBColor(0xFF, 0x8A, 0x42)
GREEN   = RGBColor(0x3D, 0xBA, 0x8A)
YELLOW  = RGBColor(0xF5, 0xC8, 0x42)

W = Inches(13.33)
H = Inches(7.5)

# "" Low-level helpers """"""""""""""""""""""""""""""""""""""""""""""""""""""""""

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, left, top, w, h, fill=SURFACE, line=BORDER, lw=Pt(0.5)):
    s = slide.shapes.add_shape(1, left, top, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.color.rgb = line; s.line.width = lw
    return s

def txb(slide, text, left, top, w, h,
        size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size    = Pt(size)
    run.font.bold    = bold
    run.font.italic  = italic
    run.font.color.rgb = color
    return box

def label(slide, text, left, top, w=Inches(8), color=ACCENT):
    txb(slide, text.upper(), left, top, w, Inches(0.25),
        size=9, bold=True, color=color)

def rule(slide, left, top, w=Inches(0.5)):
    """Thin horizontal accent line."""
    line = slide.shapes.add_shape(1, left, top, w, Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def stat_card(slide, value, lbl, left, top,
              w=Inches(3.0), h=Inches(1.2), vc=ACCENT):
    rect(slide, left, top, w, h)
    txb(slide, value, left + Inches(0.18), top + Inches(0.08),
        w - Inches(0.36), Inches(0.62), size=30, bold=True, color=vc)
    txb(slide, lbl, left + Inches(0.18), top + Inches(0.7),
        w - Inches(0.36), Inches(0.44), size=11, color=MUTED, wrap=True)

def bullet_block(slide, items, left, top, w, h,
                 size=13, color=MUTED, gap=Inches(0.34)):
    for i, item in enumerate(items):
        txb(slide, f"  -  {item}", left, top + i * gap, w, gap + Inches(0.05),
            size=size, color=color, wrap=True)

def hdr(slide, label_text, num_text):
    """Top header bar."""
    rect(slide, 0, 0, W, Inches(0.55), fill=BG, line=BORDER, lw=Pt(0.75))
    txb(slide, label_text.upper(), Inches(0.45), Inches(0.13),
        Inches(8), Inches(0.3), size=9, bold=True, color=ACCENT)
    txb(slide, "WikiMed Article Recommender", Inches(9.2), Inches(0.13),
        Inches(3.0), Inches(0.3), size=9, color=DIM, align=PP_ALIGN.RIGHT)
    txb(slide, num_text, Inches(12.3), Inches(0.13),
        Inches(0.9), Inches(0.3), size=9, color=DIM, align=PP_ALIGN.RIGHT)

# "" Slide builders """""""""""""""""""""""""""""""""""""""""""""""""""""""""""""

def slide_title(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "SC-HIA Â· Brown University", "01 / 08")

    # Background accent glow (subtle rectangle)
    glow = sl.shapes.add_shape(1, Inches(8), Inches(0.5), Inches(5.5), Inches(3.5))
    glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x00, 0x20, 0x28)
    glow.line.fill.background()

    # Content  --  bottom-anchored
    label(sl, "Research Presentation Â· July 2026", Inches(0.5), Inches(1.8))
    txb(sl, "WikiMed\nArticle\nRecommender",
        Inches(0.5), Inches(2.1), Inches(9), Inches(2.4),
        size=54, bold=True, color=TEXT)
    txb(sl, "A data-driven pipeline to identify the Wikipedia medical articles\n"
            "that need student editors most  --  ranked by impact, attention, and clinical relevance.",
        Inches(0.5), Inches(4.7), Inches(8.5), Inches(0.9),
        size=15, color=MUTED, wrap=True)

    # Meta row
    rule(sl, Inches(0.5), Inches(5.75), Inches(12.4))
    items = [
        ("Presenter", "Zach Schneider-Lynch"),
        ("Program",   "SC-HIA Â· Brown Medical School"),
        ("Mentor",    "Dr. Jeremy Warner"),
        ("Corpus",    "53,286 WikiProject Medicine articles"),
    ]
    for i, (lbl_t, val_t) in enumerate(items):
        x = Inches(0.5 + i * 3.1)
        txb(sl, lbl_t.upper(), x, Inches(5.85), Inches(2.8), Inches(0.25),
            size=8, bold=True, color=DIM)
        txb(sl, val_t, x, Inches(6.1), Inches(2.8), Inches(0.38),
            size=13, bold=(i in (0, 2)), color=MUTED if i % 2 else TEXT)


def slide_intro(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Introduction & Background", "02 / 08")

    # Left col  --  problem
    label(sl, "The Problem", Inches(0.5), Inches(0.75))
    txb(sl, "Wikipedia is where patients look first  -- \nbut most medical articles are unfinished.",
        Inches(0.5), Inches(1.0), Inches(6.0), Inches(1.1),
        size=22, bold=True, color=TEXT, wrap=True)
    rule(sl, Inches(0.5), Inches(2.2))

    problems = [
        "Wikipedia is the most-visited medical reference globally, used by ~70% of medical "
        "students and millions of patients before clinical encounters.",
        "WikiProject Medicine tracks 53,286 articles  --  yet over 60% remain at Stub or Start quality, "
        "with limited or poorly sourced content.",
        "WikiMed student editors lack a systematic way to prioritize which articles to improve. "
        "Without guidance, editing effort may miss the highest-impact targets.",
        "No existing tool combines public reach, editorial quality gaps, and clinical relevance "
        "into a single prioritization framework for volunteer editors.",
    ]
    bullet_block(sl, problems, Inches(0.5), Inches(2.35), Inches(5.9), Inches(4.0),
                 size=13, gap=Inches(1.0))

    # Right col  --  stat cards
    label(sl, "Scale", Inches(7.3), Inches(0.75))
    stat_card(sl, "53,286", "WikiProject Medicine articles tracked across all of medical Wikipedia",
              Inches(7.3), Inches(1.0), vc=ACCENT)
    stat_card(sl, "> 60%", "Articles rated Stub or Start  --  limited content, most room for improvement",
              Inches(7.3), Inches(2.3), vc=WARM)
    stat_card(sl, "~70%", "Of medical students who use Wikipedia as a clinical learning resource",
              Inches(7.3), Inches(3.6), vc=GREEN)
    stat_card(sl, "1B+", "Annual pageviews across WikiProject Medicine articles (Wikipedia Clickstream)",
              Inches(7.3), Inches(4.9), vc=YELLOW)


def slide_aims(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Research Aims & Questions", "03 / 08")

    # Aims  --  left
    label(sl, "Aims", Inches(0.5), Inches(0.75))
    aims = [
        ("A1", "Build a reproducible scoring pipeline",
         "Develop composite scores combining public reach, quality deficits, and clinical relevance "
         "to rank all 53,286 WikiProject Medicine articles by editing priority."),
        ("A2", "Create an accessible discovery tool",
         "Deploy a Streamlit dashboard enabling student editors to explore top articles by specialty, "
         "MeSH category, reading level, and real-time attention signal."),
        ("A3", "Validate against expert editorial judgment",
         "Confirm the Impact-Need Score surfaces articles independently rated as high-importance "
         "and low-quality by WikiProject Medicine volunteer editors."),
    ]
    for i, (code, title, desc) in enumerate(aims):
        y = Inches(1.05 + i * 1.82)
        rect(sl, Inches(0.5), y, Inches(5.9), Inches(1.65))
        txb(sl, code, Inches(0.65), y + Inches(0.12), Inches(0.55), Inches(0.38),
            size=11, bold=True, color=ACCENT)
        txb(sl, title, Inches(1.25), y + Inches(0.1), Inches(5.0), Inches(0.42),
            size=14, bold=True, color=TEXT)
        txb(sl, desc, Inches(1.25), y + Inches(0.54), Inches(5.0), Inches(1.0),
            size=12, color=MUTED, wrap=True)

    # RQs  --  right
    label(sl, "Research Questions", Inches(7.1), Inches(0.75))
    rqs = [
        ("RQ1", "Can a composite score combining pageviews, editorial quality, and clinical "
                "relevance produce a valid ranking of Wikipedia medical articles by editing priority?"),
        ("RQ2", "Does the Impact-Need Score identify articles rated as high-importance but "
                "low-quality by WikiProject Medicine editors  --  confirming alignment with expert judgment?"),
        ("RQ3", "What is the achievable MeSH coverage across 53,286 articles using an automated "
                "7-layer NLP cascade, and what does coverage reveal about specialty representation?"),
    ]
    for i, (code, text) in enumerate(rqs):
        y = Inches(1.05 + i * 1.82)
        # Accent border using thin rect
        accent_bar = sl.shapes.add_shape(1, Inches(7.1), y, Pt(3), Inches(1.65))
        accent_bar.fill.solid(); accent_bar.fill.fore_color.rgb = ACCENT
        accent_bar.line.fill.background()
        rect(sl, Inches(7.2), y, Inches(5.9), Inches(1.65),
             fill=RGBColor(0x05, 0x18, 0x26), line=RGBColor(0x00, 0x55, 0x65))
        txb(sl, code, Inches(7.35), y + Inches(0.12), Inches(0.6), Inches(0.38),
            size=11, bold=True, color=ACCENT)
        txb(sl, text, Inches(8.0), y + Inches(0.1), Inches(5.0), Inches(1.4),
            size=13, color=MUTED, wrap=True)


def slide_methods1(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Methods Â· Part 1 of 2  --  Data Collection", "04 / 08")

    label(sl, "Automated Pipeline", Inches(0.5), Inches(0.72))

    # Pipeline steps
    steps = [
        ("Step 1", "Article\nCollection"),
        ("Step 2", "Pageview &\nEditor Enrichment"),
        ("Step 3", "Quality &\nImportance Ratings"),
        ("Step 4", "MeSH\nAssignment"),
        ("Step 5", "Score\nComputation"),
        ("Step 6", "Dashboard\nDeployment"),
    ]
    sw = Inches(1.95)
    for i, (s, v) in enumerate(steps):
        x = Inches(0.5 + i * 2.06)
        hi = i in (1, 3)
        fill = RGBColor(0x04, 0x26, 0x35) if hi else SURFACE
        lc   = ACCENT if hi else BORDER
        rect(sl, x, Inches(0.95), sw, Inches(0.75), fill=fill, line=lc)
        txb(sl, s.upper(), x + Inches(0.1), Inches(0.98), sw - Inches(0.1), Inches(0.24),
            size=8, bold=True, color=ACCENT if hi else DIM)
        txb(sl, v, x + Inches(0.1), Inches(1.2), sw - Inches(0.1), Inches(0.45),
            size=12, color=ACCENT if hi else MUTED)
        if i < len(steps) - 1:
            txb(sl, "->", x + sw + Inches(0.02), Inches(1.05), Inches(0.12), Inches(0.35),
                size=14, color=DIM, align=PP_ALIGN.CENTER)

    # Data sources grid
    label(sl, "Data Sources", Inches(0.5), Inches(1.92))

    sources = [
        ("Wikipedia API (MediaWiki)",
         "53,286 article titles, lead sections, quality class, importance label, editor counts, inbound links, watcher counts"),
        ("Wikipedia Clickstream (May 2026)",
         "Per-article fraction of inbound clicks from external search vs. internal navigation  --  the search intent signal"),
        ("NLM MeSH 2026 (desc2026.xml)",
         "61,794 MeSH descriptors with all entry term synonyms  --  vocabulary for relevance scoring and hierarchical expansion"),
        ("PubMed API (efetch)",
         "MeSH tags from papers cited in each article; major-topic (2 pts) vs. minor-topic (1 pt) citation consensus"),
        ("Wikidata SPARQL (P486)",
         "Human-curated MeSH Descriptor ID links  --  gold-standard Layer 1 of the MeSH assignment cascade"),
        ("BioPortal Â· NCBI eSearch Â· NLM MTI",
         "NLP annotators for articles without structured data  --  Layers 5 -- 7 of cascade (NLM MTI indexes PubMed itself)"),
    ]
    cols = 2
    sw2 = Inches(6.0)
    for i, (name, desc) in enumerate(sources):
        col = i % cols
        row = i // cols
        x = Inches(0.5 + col * 6.55)
        y = Inches(2.15 + row * 1.48)
        rect(sl, x, y, sw2, Inches(1.3))
        # Dot
        dot = sl.shapes.add_shape(1, x + Inches(0.18), y + Inches(0.18),
                                   Pt(6), Pt(6))
        dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        txb(sl, name,  x + Inches(0.45), y + Inches(0.1), sw2 - Inches(0.6), Inches(0.38),
            size=13, bold=True, color=TEXT)
        txb(sl, desc,  x + Inches(0.45), y + Inches(0.48), sw2 - Inches(0.6), Inches(0.85),
            size=11, color=MUTED, wrap=True)

    txb(sl, "Reading level: Flesch-Kincaid Grade Level (textstat) from cached article lead sections (>=30 words). "
            "Rare disease flag: 13 Wikipedia categories, 1,285 articles.",
        Inches(0.5), Inches(6.98), Inches(12.0), Inches(0.3),
        size=9, color=DIM)


def slide_methods2(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Methods Â· Part 2 of 2  --  Scoring System", "05 / 08")

    label(sl, "Three Complementary Scores", Inches(0.5), Inches(0.72))

    score_data = [
        ("Impact-Need Score  Â·  0 -- 100", ACCENT,
         "Primary ranking signal",
         "Identifies articles with high public reach + low current quality.",
         "Ppv Ã—0.30  +  Nimp Ã—0.25  +  Nqd Ã—0.25  +  Pes Ã—0.10  +  Psi Ã—0.10",
         "Pageview Â· Importance label Â· Quality deficit Â· Editor scarcity Â· Search intent (Clickstream)"),
        ("Wiki Attention Score  Â·  0 -- 100", WARM,
         "Real-time momentum",
         "Measures how actively an article is being read and trending, independent of quality.",
         "Ppv Ã—0.45  +  Pvel Ã—0.20  +  Plinks Ã—0.20  +  Pwatch Ã—0.10  +  Pedit Ã—0.05",
         "Pageview Â· 3-mo traffic velocity Â· Inbound links Â· Watcher count Â· Active editors"),
        ("Medical Relevance  Â·  1 -- 10", GREEN,
         "Clinical focus signal",
         "Estimates how clinically focused an article is from text alone  --  catches biographical or historical articles under WikiProject Medicine scope.",
         "TF-IDF top-25 terms  †'  MeSH vocabulary match  Ã—  10",
         "Top 25 TF-IDF terms (unigrams + bigrams) vs. 61,794 NLM MeSH descriptors + all entry synonyms"),
    ]

    sw = Inches(4.1)
    for i, (title, col, h3, desc, formula, sub) in enumerate(score_data):
        x = Inches(0.5 + i * 4.26)
        rect(sl, x, Inches(0.95), sw, Inches(2.75))
        # Color stripe
        stripe = sl.shapes.add_shape(1, x, Inches(0.95), sw, Pt(3))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()
        txb(sl, title, x + Inches(0.18), Inches(1.02), sw - Inches(0.3), Inches(0.32),
            size=11, bold=True, color=col)
        txb(sl, h3, x + Inches(0.18), Inches(1.36), sw - Inches(0.3), Inches(0.3),
            size=14, bold=True, color=TEXT)
        txb(sl, desc, x + Inches(0.18), Inches(1.68), sw - Inches(0.3), Inches(0.52),
            size=11, color=MUTED, wrap=True)
        # Formula box
        frect = rect(sl, x + Inches(0.18), Inches(2.22), sw - Inches(0.36), Inches(0.35),
                     fill=RGBColor(0x05, 0x10, 0x18), line=col)
        txb(sl, formula, x + Inches(0.28), Inches(2.25), sw - Inches(0.55), Inches(0.28),
            size=9.5, color=col)
        txb(sl, sub, x + Inches(0.18), Inches(2.6), sw - Inches(0.3), Inches(0.3),
            size=9, color=DIM, wrap=True)

    # MeSH layers
    label(sl, "MeSH ID Assignment  --  7-Layer Cascade", Inches(0.5), Inches(3.82))

    layers = [
        ("Layer 1", "Wikidata SPARQL (P486  --  MeSH Descriptor ID)", "Exact",  GREEN),
        ("Layer 2", "Wikipedia Infobox Wikitext (MeshID= / mesh= fields)", "Exact",  GREEN),
        ("Layer 3", "MeSH Full Synonym Matching (all entry terms in desc2026.xml)", "High",   ACCENT),
        ("Layer 4", "PubMed Citation Consensus (cited paper MeSH tags, majority vote)", "High",   ACCENT),
        ("Layer 5", "BioPortal NLP Annotator (intro paragraph, MESH ontology)", "Medium", YELLOW),
        ("Layer 6", "NCBI MeSH eSearch (article title -> Metathesaurus lookup)", "Medium", YELLOW),
        ("Layer 7", "NLM Medical Text Indexer (MTI)  --  same system used to index PubMed", "Lower",  WARM),
    ]
    lh = Inches(0.46)
    for i, (lnum, lname, conf, col) in enumerate(layers):
        y = Inches(4.05 + i * lh)
        rect(sl, Inches(0.5), y, Inches(12.33), lh - Pt(2))
        txb(sl, lnum, Inches(0.65), y + Inches(0.06), Inches(0.75), Inches(0.32),
            size=10, color=DIM)
        txb(sl, lname, Inches(1.42), y + Inches(0.06), Inches(10.0), Inches(0.32),
            size=12.5, color=TEXT)
        # Confidence badge
        badge = rect(sl, Inches(11.7), y + Inches(0.08), Inches(0.95), Inches(0.28),
                     fill=RGBColor(0x05, 0x10, 0x18), line=col)
        txb(sl, conf, Inches(11.72), y + Inches(0.09), Inches(0.9), Inches(0.22),
            size=9, bold=True, color=col, align=PP_ALIGN.CENTER)


def slide_results(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Anticipated Results & Discussion", "06 / 08")

    # Left  --  charts
    label(sl, "Validation Against WikiProject Medicine Ratings", Inches(0.5), Inches(0.72))

    # Chart 1: Importance Distribution
    imp_data = ChartData()
    imp_data.categories = ["Low", "Mid", "High", "Top"]
    imp_data.add_series("All 53,286 Articles", (48, 31, 16, 5))
    imp_data.add_series("Top 100 by Score",    (3,  22, 45, 30))

    chart1 = sl.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(0.95), Inches(6.0), Inches(2.8),
        imp_data,
    ).chart
    chart1.has_title = True
    chart1.chart_title.text_frame.text = "Importance Distribution (%)"
    chart1.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    chart1.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
    chart1.has_legend = True
    from pptx.enum.chart import XL_LEGEND_POSITION
    chart1.legend.position = XL_LEGEND_POSITION.BOTTOM

    # Color series
    _color_series(chart1, 0, DIM)
    _color_series(chart1, 1, ACCENT)

    # Chart 2: Quality Distribution
    qual_data = ChartData()
    qual_data.categories = ["Stub", "Start", "C", "B", "GA/FA"]
    qual_data.add_series("All 53,286 Articles", (22, 38, 19, 13, 8))
    qual_data.add_series("Top 100 by Score",    (42, 44,  9,  4, 1))

    chart2 = sl.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5), Inches(3.85), Inches(6.0), Inches(2.8),
        qual_data,
    ).chart
    chart2.has_title = True
    chart2.chart_title.text_frame.text = "Quality Distribution (%)"
    chart2.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
    chart2.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb = MUTED
    chart2.has_legend = False

    _color_series(chart2, 0, DIM)
    _color_series(chart2, 1, ACCENT)

    # Right  --  findings
    label(sl, "Anticipated Outcomes", Inches(7.1), Inches(0.72))

    findings = [
        (ACCENT, "~82%", "MeSH descriptor coverage via 7-layer cascade  --  enabling specialty-level navigation"),
        (ACCENT, "75%",  "Of top-100 articles rated High or Top importance by WikiProject Medicine editors (vs. 21% corpus-wide)"),
        (GREEN,  "86%",  "Of top-100 articles rated Stub or Start quality  --  most room to improve (vs. 60% corpus-wide)"),
        (YELLOW, "47,319", "Articles with Flesch-Kincaid reading level computed  --  mean grade 14.4"),
    ]
    for i, (col, val, txt) in enumerate(findings):
        y = Inches(0.95 + i * 1.37)
        # Accent left bar
        bar = sl.shapes.add_shape(1, Inches(7.1), y, Pt(3.5), Inches(1.2))
        bar.fill.solid(); bar.fill.fore_color.rgb = col
        bar.line.fill.background()
        rect(sl, Inches(7.25), y, Inches(5.8), Inches(1.2),
             fill=RGBColor(0x05, 0x10, 0x18), line=RGBColor(0x0A, 0x28, 0x38))
        txb(sl, val, Inches(7.42), y + Inches(0.06), Inches(2.5), Inches(0.6),
            size=28, bold=True, color=col)
        txb(sl, txt, Inches(7.42), y + Inches(0.65), Inches(5.5), Inches(0.52),
            size=12, color=MUTED, wrap=True)

    # Demo callout
    cbox = rect(sl, Inches(7.1), Inches(6.46), Inches(5.95), Inches(0.72),
                fill=RGBColor(0x1A, 0x0C, 0x02),
                line=RGBColor(0x55, 0x30, 0x10))
    txb(sl, "->  Live dashboard demonstration to follow",
        Inches(7.28), Inches(6.54), Inches(5.6), Inches(0.3),
        size=13, bold=True, color=WARM)


def slide_validation(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "Validation Plan", "07 / 08")
    label(sl, "Three-Stage Validation Strategy", Inches(0.5), Inches(0.72))

    # "" Layout constants """""""""""""""""""""""""""""""""""""""""""""""""""""""
    col_w  = Inches(4.11)
    col_xs = [Inches(0.5), Inches(4.61), Inches(8.72)]
    nd     = Inches(0.44)           # node diameter
    node_y = Inches(1.0)
    node_cy = node_y + nd / 2       # node vertical center (for connector)
    card_y = node_y + nd + Inches(0.18)
    card_h = Inches(5.55)

    # "" Connector line across all three nodes """"""""""""""""""""""""""""""""""
    line_l = col_xs[0] + col_w / 2
    line_r = col_xs[2] + col_w / 2
    conn = sl.shapes.add_shape(1, line_l, node_cy - Pt(1), line_r - line_l, Pt(2))
    conn.fill.solid(); conn.fill.fore_color.rgb = DIM
    conn.line.fill.background()

    # "" Stage definitions """"""""""""""""""""""""""""""""""""""""""""""""""""""
    stages = [
        {
            "col":    GREEN,
            "badge":  "COMPLETED",
            "label":  "Stage 1",
            "name":   "Formula Calibration Check",
            "date":   "July 2026",
            "node_fill": GREEN,
            "node_txt":  "checkmark",
            "content": "items",
            "items": [
                "84% of top-100 ranked articles rated Stub or Start quality -- formula correctly surfaces improvement-eligible articles",
                "14x enrichment of High/Top importance articles vs. full dataset baseline (30% vs. 2.2%)",
                "Pipeline confirmed working as designed -- note: quality & importance are direct inputs (50% of formula); independent validity tested in Stages 2-3",
            ],
        },
        {
            "col":    ACCENT,
            "badge":  "IN PROGRESS",
            "label":  "Stage 2",
            "name":   "Expert & User Review",
            "date":   "July to Early August 2026",
            "node_fill": ACCENT,
            "node_txt":  None,
            "content": "parts",
            "parts": [
                ("PART 1  --  Science Librarians",
                 "Expert review of MeSH assignment methodology, search accuracy, and scoring logic validity"),
                ("PART 2  --  Focus Group",
                 "Dashboard design feedback and usability testing from a student editor perspective"),
            ],
        },
        {
            "col":    WARM,
            "badge":  "UPCOMING",
            "label":  "Stage 3",
            "name":   "WikiMed Student Deployment",
            "date":   "Mid-August 2026",
            "node_fill": None,
            "node_txt":  None,
            "content": "schools",
            "schools": [
                "Brown University",
                "Vanderbilt University",
                "+ 4 additional medical schools (potential)",
            ],
            "footnote": "Real-world usage testing; editing outcomes tracked pre- and post-dashboard introduction",
        },
    ]

    for x, stage in zip(col_xs, stages):
        col = stage["col"]

        # "" Node circle """"""""""""""""""""""""""""""""""""""""""""""""""""""""
        node = sl.shapes.add_shape(9, x + col_w / 2 - nd / 2, node_y, nd, nd)
        if stage["node_fill"]:
            node.fill.solid(); node.fill.fore_color.rgb = stage["node_fill"]
            node.line.fill.background()
        else:
            node.fill.background()
            node.line.color.rgb = col; node.line.width = Pt(2)

        if stage["node_txt"] == "checkmark":
            tf = node.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            run = p.add_run(); run.text = "+"     # closest to checkmark in basic fonts
            run.font.size = Pt(18); run.font.bold = True
            run.font.color.rgb = BG

        # "" Card """""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""
        rect(sl, x, card_y, col_w, card_h)
        stripe = sl.shapes.add_shape(1, x, card_y, col_w, Pt(3))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col
        stripe.line.fill.background()

        cy = card_y + Inches(0.12)

        # Badge
        txb(sl, stage["badge"], x + Inches(0.18), cy,
            col_w - Inches(0.3), Inches(0.25), size=8, bold=True, color=col)

        # Stage label + name + date
        txb(sl, stage["label"].upper(), x + Inches(0.18), cy + Inches(0.28),
            col_w - Inches(0.3), Inches(0.2), size=8, color=DIM)
        txb(sl, stage["name"], x + Inches(0.18), cy + Inches(0.5),
            col_w - Inches(0.3), Inches(0.42), size=13, bold=True, color=TEXT)
        txb(sl, stage["date"], x + Inches(0.18), cy + Inches(0.95),
            col_w - Inches(0.3), Inches(0.25), size=10.5, color=MUTED)

        # Divider
        div = sl.shapes.add_shape(1, x + Inches(0.18), cy + Inches(1.28),
                                   col_w - Inches(0.36), Pt(0.75))
        div.fill.solid(); div.fill.fore_color.rgb = BORDER
        div.line.fill.background()

        content_y = cy + Inches(1.42)

        # "" Stage-specific content """""""""""""""""""""""""""""""""""""""""""""
        if stage["content"] == "items":
            ih = Inches(1.18)
            for j, item in enumerate(stage["items"]):
                iy = content_y + j * (ih + Inches(0.08))
                r = rect(sl, x + Inches(0.18), iy, col_w - Inches(0.36), ih,
                         fill=RGBColor(0x04, 0x18, 0x14), line=col)
                # Check dot
                dot = sl.shapes.add_shape(1, x + Inches(0.3), iy + Inches(0.14),
                                           Pt(5), Pt(5))
                dot.fill.solid(); dot.fill.fore_color.rgb = col
                dot.line.fill.background()
                txb(sl, item, x + Inches(0.52), iy + Inches(0.08),
                    col_w - Inches(0.75), ih - Inches(0.16),
                    size=10.5, color=MUTED, wrap=True)

        elif stage["content"] == "parts":
            ph = Inches(1.58)
            for j, (plabel, ptext) in enumerate(stage["parts"]):
                py = content_y + j * (ph + Inches(0.1))
                rect(sl, x + Inches(0.18), py, col_w - Inches(0.36), ph,
                     fill=RGBColor(0x03, 0x18, 0x22), line=col)
                txb(sl, plabel, x + Inches(0.28), py + Inches(0.1),
                    col_w - Inches(0.55), Inches(0.25), size=8, bold=True, color=col)
                txb(sl, ptext, x + Inches(0.28), py + Inches(0.38),
                    col_w - Inches(0.55), ph - Inches(0.48),
                    size=11, color=MUTED, wrap=True)

        elif stage["content"] == "schools":
            sh = Inches(0.52)
            for j, school in enumerate(stage["schools"]):
                is_note = (j == 2)
                txb(sl, school,
                    x + Inches(0.18), content_y + j * sh,
                    col_w - Inches(0.36), sh - Inches(0.04),
                    size=12, italic=is_note,
                    color=MUTED if is_note else TEXT)
                if not is_note:
                    sep = sl.shapes.add_shape(
                        1, x + Inches(0.18), content_y + j * sh + sh - Pt(1),
                        col_w - Inches(0.36), Pt(0.75))
                    sep.fill.solid(); sep.fill.fore_color.rgb = BORDER
                    sep.line.fill.background()

            fn_y = content_y + len(stage["schools"]) * sh + Inches(0.15)
            txb(sl, stage["footnote"],
                x + Inches(0.18), fn_y, col_w - Inches(0.36), Inches(0.7),
                size=9, color=DIM, wrap=True)


def slide_refs(prs):
    sl = blank_slide(prs); set_bg(sl)
    hdr(sl, "References", "08 / 08")

    label(sl, "Key Sources", Inches(0.5), Inches(0.72))

    refs = [
        ("1.", "Heilman JM, et al.", "Wikipedia: a key tool for global public health promotion.",
         "J Med Internet Res. 2011;13(1):e14."),
        ("2.", "Laurent MR, Vickers TJ.", "Seeking health information online: does Wikipedia matter?",
         "J Am Med Inform Assoc. 2009;16(4):471 -- 479."),
        ("3.", "Reavley NJ, et al.", "Quality of information sources about mental disorders: a comparison of Wikipedia with controlled web and print sources.",
         "Psychol Med. 2012;42(8):1753 -- 1762."),
        ("4.", "National Library of Medicine.", "Medical Subject Headings (MeSH) 2026 Descriptor File.",
         "Bethesda, MD: NLM; 2026."),
        ("5.", "Noy NF, et al.", "BioPortal: ontologies and integrated data resources at the click of a mouse.",
         "Nucleic Acids Res. 2009;37(suppl 2):W170 -- W173."),
        ("6.", "Wikimedia Foundation.", "Wikipedia Clickstream.",
         "Wikimedia Research Data Repository; May 2026."),
        ("7.", "Flesch R.", "A new readability yardstick.",
         "J Appl Psychol. 1948;32(3):221 -- 233."),
        ("8.", "Wikipedia.", "WikiProject Medicine.",
         "en.wikipedia.org/wiki/Wikipedia:WikiProject_Medicine"),
    ]

    cols = 2
    rh = Inches(0.72)
    for i, (num, author, title, source) in enumerate(refs):
        col = i % cols
        row = i // cols
        x = Inches(0.5  + col * 6.55)
        y = Inches(1.0 + row * rh)
        rect(sl, x, y, Inches(6.1), rh - Pt(4), fill=BG, line=BORDER)
        # Left accent
        acc = sl.shapes.add_shape(1, x, y, Pt(2.5), rh - Pt(4))
        acc.fill.solid(); acc.fill.fore_color.rgb = BORDER
        acc.line.fill.background()
        txb(sl, f"{num} {author} {title} {source}",
            x + Inches(0.18), y + Inches(0.08), Inches(5.8), rh - Inches(0.1),
            size=12, color=MUTED, wrap=True)

    txb(sl, "Dashboard source code, scoring pipeline, and enrichment scripts available upon request. "
            "Data: Wikipedia API Â· Wikidata SPARQL Â· NLM MeSH 2026 Â· Wikipedia Clickstream Â· "
            "PubMed API Â· BioPortal Â· NCBI eSearch Â· NLM MTI.",
        Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
        size=9, color=DIM)


def _color_series(chart, idx, rgb):
    """Set a solid fill color for a chart series by index."""
    try:
        ser = chart.series[idx]
        fill = ser.format.fill
        fill.solid()
        fill.fore_color.rgb = rgb
    except Exception:
        pass


# "" Main """""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""""

def main():
    prs = new_prs()

    print("Building slides...")
    slide_title(prs);      print("  1/8 Title")
    slide_intro(prs);      print("  2/8 Introduction")
    slide_aims(prs);       print("  3/8 Research Aims")
    slide_methods1(prs);   print("  4/8 Methods 1")
    slide_methods2(prs);   print("  5/8 Methods 2")
    slide_results(prs);    print("  6/8 Results")
    slide_validation(prs); print("  7/8 Validation Plan")
    slide_refs(prs);       print("  8/8 References")

    out = "WikiMed_Presentation.pptx"
    prs.save(out)
    print(f"\nSaved: {out}")
    print("\nTo open in Google Slides:")
    print("  1. Go to drive.google.com")
    print("  2. Click New > File upload")
    print(f"  3. Upload {out}")
    print("  4. Double-click the file -- it opens directly in Google Slides")


if __name__ == "__main__":
    main()

